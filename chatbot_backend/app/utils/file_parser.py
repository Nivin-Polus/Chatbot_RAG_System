# app/utils/file_parser.py

from io import BytesIO
from typing import List
from math import ceil

# PDF
from pypdf import PdfReader

# DOCX
from docx import Document

# PPTX
from pptx import Presentation

# Excel
import pandas as pd

# Chunk size (number of words per chunk)
CHUNK_SIZE = 200


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE) -> List[str]:
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunks.append(" ".join(words[i:i + chunk_size]))
    return chunks


def parse_file(filename: str, content: bytes, chunk_size: int = CHUNK_SIZE) -> List[str]:
    """
    Extract text from file and split into chunks.
    Supports PDF, DOCX, PPTX, XLS/XLSX, TXT, CSV.
    """
    import logging
    logger = logging.getLogger("file_parser")
    
    ext = filename.split('.')[-1].lower()
    text = ""
    
    logger.info(f"[FILE PARSER] Processing file: {filename}, extension: {ext}, size: {len(content)} bytes")

    try:
        if ext == "pdf":
            logger.info(f"[FILE PARSER] Processing PDF file")
            reader = PdfReader(BytesIO(content))
            page_count = 0
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                    page_count += 1
            logger.info(f"[FILE PARSER] PDF: Extracted text from {page_count} pages")

        elif ext == "docx":
            logger.info(f"[FILE PARSER] Processing DOCX file")
            doc = Document(BytesIO(content))
            para_count = 0
            for para in doc.paragraphs:
                if para.text:
                    text += para.text + "\n"
                    para_count += 1
            logger.info(f"[FILE PARSER] DOCX: Extracted {para_count} paragraphs")

        elif ext == "pptx":
            logger.info(f"[FILE PARSER] Processing PPTX file")
            prs = Presentation(BytesIO(content))
            slide_count = 0
            shape_count = 0
            for slide in prs.slides:
                slide_count += 1
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                        shape_count += 1
            logger.info(f"[FILE PARSER] PPTX: Processed {slide_count} slides, {shape_count} text shapes")

        elif ext in ["xls", "xlsx"]:
            logger.info(f"[FILE PARSER] Processing Excel file")
            xls = pd.ExcelFile(BytesIO(content))
            sheet_count = len(xls.sheet_names)
            row_count = 0
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                for row in df.itertuples(index=False):
                    row_text = " ".join([str(cell) for cell in row if pd.notna(cell)])
                    if row_text.strip():
                        text += row_text + "\n"
                        row_count += 1
            logger.info(f"[FILE PARSER] Excel: Processed {sheet_count} sheets, {row_count} rows")

        elif ext == "txt":
            logger.info(f"[FILE PARSER] Processing TXT file")
            try:
                text = content.decode("utf-8")
                logger.info(f"[FILE PARSER] TXT: Decoded with UTF-8")
            except UnicodeDecodeError:
                text = content.decode("latin-1")
                logger.info(f"[FILE PARSER] TXT: Decoded with Latin-1")

        elif ext == "csv":
            logger.info(f"[FILE PARSER] Processing CSV file")
            try:
                # First try basic text decoding to ensure file is readable
                try:
                    decoded_content = content.decode("utf-8")
                except UnicodeDecodeError:
                    try:
                        decoded_content = content.decode("latin-1")
                    except UnicodeDecodeError:
                        decoded_content = content.decode("utf-8", errors="ignore")
                
                # Try to read CSV with pandas
                try:
                    df = pd.read_csv(BytesIO(content))
                    row_count = len(df)
                    col_count = len(df.columns)
                    
                    # Convert CSV to readable text format
                    text_lines = []
                    
                    # Add header row
                    header = " | ".join(df.columns.astype(str))
                    text_lines.append(f"Headers: {header}")
                    text_lines.append("-" * len(header))
                    
                    # Add data rows (limit to prevent huge files)
                    max_rows = min(1000, row_count)  # Limit to 1000 rows for performance
                    for index, row in df.head(max_rows).iterrows():
                        row_text = " | ".join([str(cell) if pd.notna(cell) else "" for cell in row])
                        if row_text.strip():
                            text_lines.append(f"Row {index + 1}: {row_text}")
                    
                    if row_count > max_rows:
                        text_lines.append(f"... and {row_count - max_rows} more rows")
                    
                    text = "\n".join(text_lines)
                    logger.info(f"[FILE PARSER] CSV: Processed {min(max_rows, row_count)} of {row_count} rows, {col_count} columns")
                    
                except Exception as pandas_error:
                    logger.warning(f"[FILE PARSER] CSV pandas parsing failed: {pandas_error}")
                    # Fallback: treat as plain text
                    text = decoded_content
                    logger.info(f"[FILE PARSER] CSV: Using fallback text parsing")
                
            except Exception as csv_error:
                logger.error(f"[FILE PARSER] CSV processing failed: {csv_error}")
                # Last resort: return minimal text
                text = f"CSV file content (parsing failed): {filename}"
                logger.info(f"[FILE PARSER] CSV: Using minimal fallback")

        else:
            logger.error(f"[FILE PARSER] Unsupported file type: {ext}")
            raise ValueError(f"Unsupported file type: {ext}")

    except Exception as e:
        logger.error(f"[FILE PARSER] Error parsing file {filename}: {str(e)}")
        raise ValueError(f"Error parsing file {filename}: {str(e)}")

    # Log extracted text info
    logger.info(f"[FILE PARSER] Extracted text length: {len(text)} characters")
    
    # Split text into chunks
    chunks = chunk_text(text, chunk_size)
    logger.info(f"[FILE PARSER] Created {len(chunks)} chunks from {filename}")
    
    return chunks
